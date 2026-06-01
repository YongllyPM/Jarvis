"""file_processor.py — Processes uploaded files (PDFs, images, documents, etc.)."""

import os, json, csv, io, traceback, zipfile, shutil
from pathlib import Path
from datetime import datetime

def file_processor(parameters: dict, player=None, speak=None) -> str:
    action = (parameters.get("action") or "info").lower()
    file_path_str = parameters.get("file_path", "")
    instruction = parameters.get("instruction", "")
    dest_format = parameters.get("format", "")
    width = parameters.get("width")
    height = parameters.get("height")
    scale = parameters.get("scale")
    quality = parameters.get("quality")
    start = parameters.get("start")
    end = parameters.get("end")
    timestamp = parameters.get("timestamp")
    column = parameters.get("column")
    value = parameters.get("value")
    condition = parameters.get("condition")
    ascending = parameters.get("ascending")
    save = parameters.get("save", True)
    destination = parameters.get("destination", "")

    if not file_path_str:
        return "Error: No se especificó la ruta del archivo."

    path = Path(file_path_str)
    if not path.exists():
        return f"Error: El archivo '{file_path_str}' no existe."

    ext = path.suffix.lower()

    try:
        if ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff"):
            return _process_image(path, action, instruction, dest_format, width, height, scale, quality, save, player, speak)
        elif ext == ".pdf":
            return _process_pdf(path, action, instruction, dest_format, save, player, speak)
        elif ext in (".docx", ".doc"):
            return _process_docx(path, action, instruction, dest_format, save, player, speak)
        elif ext in (".txt", ".md", ".log"):
            return _process_text(path, action, instruction, dest_format, save)
        elif ext in (".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp", ".h", ".hpp",
                     ".cs", ".rb", ".go", ".rs", ".swift", ".kt", ".scala", ".php", ".r",
                     ".sql", ".sh", ".bat", ".ps1", ".yml", ".yaml", ".toml", ".ini", ".cfg",
                     ".html", ".css", ".scss", ".less", ".vue", ".svelte"):
            return _process_code(path, action, instruction, save)
        elif ext in (".csv", ".tsv"):
            return _process_csv(path, action, instruction, column, value, condition, ascending, dest_format, save)
        elif ext in (".json", ".xml"):
            return _process_data(path, action, instruction, dest_format, save)
        elif ext in (".mp3", ".wav", ".flac", ".ogg", ".m4a", ".aac", ".wma"):
            return _process_audio(path, action, instruction, start, end, dest_format, save)
        elif ext in (".mp4", ".avi", ".mkv", ".mov", ".wmv", ".webm", ".flv"):
            return _process_video(path, action, instruction, start, end, timestamp, dest_format, quality, save)
        elif ext in (".zip", ".rar", ".7z", ".tar", ".gz"):
            return _process_archive(path, action, destination, save)
        elif ext in (".pptx", ".ppt"):
            return _process_pptx(path, action, instruction, save)
        else:
            return _process_unknown(path, action, instruction)

    except Exception as e:
        return f"Error al procesar '{path.name}': {str(e)}\n{traceback.format_exc()}"


def _get_size_str(size_bytes: int) -> str:
    for unit in ("B", "KB", "MB", "GB"):
        if size_bytes < 1024:
            return f"{size_bytes:.1f} {unit}"
        size_bytes /= 1024
    return f"{size_bytes:.1f} TB"


def _info_dict(path: Path) -> dict:
    return {
        "nombre": path.name,
        "tamaño": _get_size_str(os.path.getsize(path)),
        "extensión": path.suffix.lower(),
        "modificado": datetime.fromtimestamp(os.path.getmtime(path)).strftime("%Y-%m-%d %H:%M")
    }


# ──────────────────────────── PDF ────────────────────────────

def _process_pdf(path, action, instruction, dest_format, save, player, speak):
    try:
        from pypdf import PdfReader
    except ImportError:
        return "Error: La librería pypdf no está instalada. No se puede procesar el PDF."

    reader = PdfReader(str(path))
    num_pages = len(reader.pages)
    text = ""
    for i, page in enumerate(reader.pages):
        page_text = page.extract_text() or ""
        text += f"\n--- Página {i+1} ---\n{page_text}"

    if action == "info":
        info = _info_dict(path)
        info["páginas"] = num_pages
        return json.dumps(info, ensure_ascii=False)

    if action == "extract_text" or action == "summarize":
        if not text.strip():
            return "No se pudo extraer texto del PDF (puede ser un PDF escaneado)."
        truncated = text[:8000]
        suffix = f"\n\n[... {len(text)} caracteres totales, mostrando primeros 8000]" if len(text) > 8000 else ""
        return f"Texto extraído de '{path.name}' ({num_pages} págs):\n\n{truncated}{suffix}"

    return f"'{path.name}' — {num_pages} páginas, {_get_size_str(os.path.getsize(path))}. Usa 'extract_text' para leerlo."


# ──────────────────────────── IMAGE ────────────────────────────

def _process_image(path, action, instruction, dest_format, width, height, scale, quality, save, player, speak):
    from PIL import Image

    img = Image.open(str(path))

    if action == "info":
        info = _info_dict(path)
        info["dimensiones"] = f"{img.width}x{img.height}"
        info["formato"] = img.format
        info["modo"] = img.mode
        return json.dumps(info, ensure_ascii=False)

    if action in ("ocr", "describe", "recognize"):
        try:
            import pytesseract
            ocr_result = pytesseract.image_to_string(img, lang="spa+eng")
            if ocr_result.strip():
                return f"Texto detectado en la imagen:\n\n{ocr_result.strip()[:4000]}"
            return "No se detectó texto en la imagen."
        except Exception as e:
            return f"Error al aplicar OCR: {e}"

    if action == "resize":
        w = width or 800
        h = height or 600
        img = img.resize((w, h), Image.LANCZOS)
        if save:
            out_path = path.parent / f"{path.stem}_resized{path.suffix}"
            img.save(str(out_path))
            return f"Imagen redimensionada a {w}x{h} → '{out_path.name}'"
        return f"Imagen redimensionada a {w}x{h}."

    if action == "compress":
        q = quality or 75
        if save:
            out_path = path.parent / f"{path.stem}_compressed{path.suffix}"
            img.save(str(out_path), quality=q, optimize=True)
            sz_before = _get_size_str(os.path.getsize(path))
            sz_after = _get_size_str(os.path.getsize(out_path))
            return f"Imagen comprimida (calidad {q}%): {sz_before} → {sz_after} → '{out_path.name}'"
        return f"Imagen comprimida al {q}% de calidad."

    if action == "convert" and dest_format:
        fmt = dest_format.upper()
        if fmt in ("JPG", "JPEG"):
            fmt = "JPEG"
        out_path = path.parent / f"{path.stem}.{dest_format.lower()}"
        img.convert("RGB").save(str(out_path), fmt)
        return f"Imagen convertida a {fmt} → '{out_path.name}'"

    return f"'{path.name}' — {img.width}x{img.height}, {img.format}. Acciones: ocr, resize, compress, convert."


# ──────────────────────────── DOCX ────────────────────────────

def _process_docx(path, action, instruction, dest_format, save, player, speak):
    try:
        from docx import Document
    except ImportError:
        return "Error: python-docx no está instalado."

    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    full_text = "\n".join(paragraphs)

    if action == "info":
        info = _info_dict(path)
        info["párrafos"] = len(paragraphs)
        return json.dumps(info, ensure_ascii=False)

    if action == "extract_text" or action == "summarize":
        truncated = full_text[:6000]
        suffix = f"\n\n[... {len(full_text)} caracteres totales]" if len(full_text) > 6000 else ""
        return f"Texto de '{path.name}':\n\n{truncated}{suffix}"

    return f"'{path.name}' — {len(paragraphs)} párrafos."


# ──────────────────────────── TEXT ────────────────────────────

def _process_text(path, action, instruction, dest_format, save):
    text = path.read_text(encoding="utf-8", errors="replace")

    if action == "info":
        lines = text.count("\n") + 1
        words = len(text.split())
        info = _info_dict(path)
        info["líneas"] = lines
        info["palabras"] = words
        info["caracteres"] = len(text)
        return json.dumps(info, ensure_ascii=False)

    if action == "summarize":
        return f"Resumen de '{path.name}':\n\n{text[:3000]}"

    if action == "word_count":
        return f"Palabras: {len(text.split())}, Caracteres: {len(text)}, Líneas: {text.count(chr(10))+1}"

    if action == "to_bullet":
        lines = text.strip().split("\n")
        bulleted = "\n".join(f"• {l}" for l in lines if l.strip())
        return f"'{path.name}' en viñetas:\n\n{bulleted[:4000]}"

    return f"Texto de '{path.name}':\n\n{text[:3000]}"


# ──────────────────────────── CODE ────────────────────────────

def _process_code(path, action, instruction, save):
    code = path.read_text(encoding="utf-8", errors="replace")
    lang = path.suffix[1:] or "texto"

    if action == "info":
        lines = code.count("\n") + 1
        info = _info_dict(path)
        info["lenguaje"] = lang
        info["líneas"] = lines
        return json.dumps(info, ensure_ascii=False)

    if action == "explain":
        return f"Explicación del código '{path.name}' ({lang}, {code.count(chr(10))+1} líneas):\n\n```{lang}\n{code[:4000]}\n```"

    if action == "review" or action == "fix" or action == "optimize":
        return f"Código '{path.name}' ({lang}):\n\n```{lang}\n{code[:5000]}\n```\n\n(Pasa el código al asistente principal para revisión detallada.)"

    return f"'{path.name}' — {lang}, {code.count(chr(10))+1} líneas."


# ──────────────────────────── CSV ────────────────────────────

def _process_csv(path, action, instruction, column, value, condition, ascending, dest_format, save):
    import csv
    rows = []
    with open(str(path), newline="", encoding="utf-8-sig") as f:
        reader = csv.DictReader(f)
        headers = reader.fieldnames or []
        for r in reader:
            rows.append(r)

    if action == "info" or action == "analyze":
        info = _info_dict(path)
        info["columnas"] = headers
        info["filas"] = len(rows)
        preview = rows[:5] if rows else []
        return json.dumps({"info": info, "preview": preview}, ensure_ascii=False, indent=2)

    if action == "stats":
        return f"'{path.name}': {len(headers)} columnas, {len(rows)} filas.\nColumnas: {', '.join(headers)}"

    if action == "filter" and column:
        filtered = []
        for r in rows:
            v = r.get(column, "")
            if condition == "contains":
                if value and value.lower() in v.lower():
                    filtered.append(r)
            elif condition == "gt":
                try: filtered.append(r) if float(v) > float(value) else None
                except: pass
            elif condition == "lt":
                try: filtered.append(r) if float(v) < float(value) else None
                except: pass
            else:
                if v == value:
                    filtered.append(r)
        return json.dumps({"filtrados": len(filtered), "filas": filtered[:20]}, ensure_ascii=False, indent=2)

    if action == "sort" and column:
        reverse = not (ascending if ascending is not None else True)
        rows.sort(key=lambda r: (r.get(column) or "").lower(), reverse=reverse)
        if save:
            out = path.parent / f"{path.stem}_sorted{path.suffix}"
            with open(str(out), "w", newline="", encoding="utf-8") as f:
                w = csv.DictWriter(f, fieldnames=headers)
                w.writeheader()
                w.writerows(rows)
            return f"Ordenado por '{column}' → '{out.name}'"
        return f"Ordenado por '{column}' ({len(rows)} filas)."

    return f"'{path.name}' — {len(headers)} columnas, {len(rows)} filas."


# ──────────────────────────── JSON/XML ────────────────────────────

def _process_data(path, action, instruction, dest_format, save):
    raw = path.read_text(encoding="utf-8", errors="replace")

    if action == "validate":
        try:
            json.loads(raw)
            return f"'{path.name}' — JSON válido."
        except json.JSONDecodeError as e:
            return f"'{path.name}' — JSON inválido: {e}"

    if action == "format":
        try:
            data = json.loads(raw)
            formatted = json.dumps(data, ensure_ascii=False, indent=2)
            if save:
                path.write_text(formatted, encoding="utf-8")
                return f"'{path.name}' formateado."
            return formatted[:4000]
        except json.JSONDecodeError as e:
            return f"Error al formatear: {e}"

    if action == "analyze":
        try:
            data = json.loads(raw)
            return json.dumps({"archivo": path.name, "tipo": type(data).__name__,
                               "tamaño": _get_size_str(os.path.getsize(path))}, ensure_ascii=False)
        except json.JSONDecodeError:
            return f"'{path.name}' — {_get_size_str(os.path.getsize(path))} (no es JSON válido)"

    return f"'{path.name}' — {_get_size_str(os.path.getsize(path))}."


# ──────────────────────────── AUDIO ────────────────────────────

def _process_audio(path, action, instruction, start, end, dest_format, save):
    if action == "info":
        info = _info_dict(path)
        return json.dumps(info, ensure_ascii=False)

    return f"'{path.name}' — {_get_size_str(os.path.getsize(path))}. Para transcribir, usa el asistente principal."


# ──────────────────────────── VIDEO ────────────────────────────

def _process_video(path, action, instruction, start, end, timestamp, dest_format, quality, save):
    if action == "info":
        info = _info_dict(path)
        return json.dumps(info, ensure_ascii=False)

    return f"'{path.name}' — {_get_size_str(os.path.getsize(path))}. Usa el asistente principal para procesar este video."


# ──────────────────────────── ARCHIVE ────────────────────────────

def _process_archive(path, action, destination, save):
    if action == "info" or action == "list":
        try:
            import zipfile
            with zipfile.ZipFile(str(path), "r") as z:
                names = z.namelist()
            return f"'{path.name}' contiene {len(names)} archivos:\n" + "\n".join(names[:50])
        except zipfile.BadZipFile:
            return f"'{path.name}' no es un zip válido o no soportado."

    if action == "extract":
        dest = Path(destination) if destination else path.parent / path.stem
        dest.mkdir(parents=True, exist_ok=True)
        try:
            import zipfile
            with zipfile.ZipFile(str(path), "r") as z:
                z.extractall(str(dest))
            return f"Extraído → '{dest}/'"
        except zipfile.BadZipFile:
            return f"No se pudo extraer '{path.name}'."

    return f"'{path.name}' — {_get_size_str(os.path.getsize(path))}."


# ──────────────────────────── PPTX ────────────────────────────

def _process_pptx(path, action, instruction, save):
    try:
        from pptx import Presentation
    except ImportError:
        return "python-pptx no está instalado."

    prs = Presentation(str(path))
    slides_text = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides_text.append(f"--- Slide {i+1} ---\n" + "\n".join(texts))

    if action == "info":
        info = _info_dict(path)
        info["diapositivas"] = len(prs.slides)
        return json.dumps(info, ensure_ascii=False)

    if action == "extract_text" or action == "summarize":
        full = "\n\n".join(slides_text)
        return f"Contenido de '{path.name}' ({len(prs.slides)} diapositivas):\n\n{full[:6000]}"

    return f"'{path.name}' — {len(prs.slides)} diapositivas."


# ──────────────────────────── UNKNOWN ────────────────────────────

def _process_unknown(path, action, instruction):
    return f"'{path.name}' — {_get_size_str(os.path.getsize(path))}. Tipo '{path.suffix}' no reconocido. Usa el asistente principal para más ayuda."
